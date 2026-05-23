import collections 
import collections.abc
from pptx import Presentation
import os

prs = Presentation('Exposición_Formato.pptx')
print(f"Total slides: {len(prs.slides)}")

for i, slide in enumerate(prs.slides):
    print(f"\n--- Slide {i+1} ---")
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.text.strip():
                    print(paragraph.text)
